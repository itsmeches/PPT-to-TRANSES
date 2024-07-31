from flask import Flask, request, send_file, render_template
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from PIL import Image
import io
import os

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def sanitize_text(text):
    """Remove non-XML compatible characters from text."""
    return ''.join(char for char in text if char.isprintable())

def is_important_image(shape):
    """Determine if an image is important based on size or position."""
    # Example filter: Only include images larger than 50x50 pixels
    if shape.width > 500000 and shape.height > 500000:  # Adjust size thresholds as needed
        return True
    return False

def set_two_column_layout(doc):
    """Set the document to two-column layout."""
    section = doc.sections[-1]
    sectPr = section._sectPr
    cols = sectPr.xpath('./w:cols')[0]
    cols.set(qn('w:num'), '2')

def pptx_to_docx(pptx_file):
    prs = Presentation(pptx_file)
    doc = Document()

    # Set document to two-column layout
    set_two_column_layout(doc)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                sanitized_text = sanitize_text(shape.text)
                if sanitized_text:
                    p = doc.add_paragraph()
                    p.add_run(sanitized_text)
                    # Add a line break after each paragraph to ensure proper separation
                    p.add_run().add_break()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER if shape.text_frame.paragraphs[0].alignment == 1 else WD_ALIGN_PARAGRAPH.LEFT
            elif shape.shape_type == 13 and hasattr(shape, "image"):  # picture
                if is_important_image(shape):
                    try:
                        image = shape.image
                        image_bytes = io.BytesIO(image.blob)
                        doc.add_picture(image_bytes, width=Inches(4))  # Adjust size as needed
                    except AttributeError:
                        continue
    
    output_path = os.path.join(UPLOAD_FOLDER, "output.docx")
    doc.save(output_path)
    return output_path

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'No file part', 400
    file = request.files['file']
    if file.filename == '':
        return 'No selected file', 400
    if file and file.filename.endswith('.pptx'):
        pptx_file = file
        docx_path = pptx_to_docx(pptx_file)
        return send_file(docx_path, as_attachment=True)
    return 'Invalid file type', 400

if __name__ == '__main__':
    app.run(debug=True)
